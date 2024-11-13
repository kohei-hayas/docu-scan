import fitz # PyMuPDF for PDF text extraction
from pptx import Presentation
from docx import Document

def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF using PyMuPDF."""
    doc = fitz.open(pdf_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_pptx(pptx_path):
    """Extract text from a PowerPoint file."""
    powerpoint = Presentation(pptx_path)
    text = ""
    for slide in powerpoint.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_docx(docx_path):
    """Extract text from a Word document."""
    doc = Document(docx_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text