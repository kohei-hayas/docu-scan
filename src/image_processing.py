import fitz  # PyMuPDF for extracting embedded images in PDFs
import pytesseract
from pptx import Presentation
from docx import Document
from src.utils import preprocess_image, convert_bytes_to_pil_image  # Import utility functions

def extract_image_text_pdf(pdf_path):
    """Extract embedded images from a PDF and apply OCR with preprocessing."""
    doc = fitz.open(pdf_path)
    text = ""

    for page_index in range(len(doc)):
        # Load each page in the PDF
        page = doc.load_page(page_index)
        # Get all images on the page
        images = page.get_images(full=True)
        
        for img in images:
            xref = img[0]  # xref is the reference number of the image in the PDF
            base_image = doc.extract_image(xref)  # Extract image using the reference number
            image_bytes = base_image["image"]  # Get the raw image bytes

            # Convert raw bytes to a PIL Image using the utility function
            pil_image = convert_bytes_to_pil_image(image_bytes)

            # Preprocess the image before OCR
            preprocessed_image = preprocess_image(pil_image)

            # Apply OCR to the preprocessed image
            text += pytesseract.image_to_string(preprocessed_image) + ' '
    
    return text


def extract_image_text_pptx(pptx_path):
    """Extract images from a PowerPoint file and apply OCR with preprocessing."""
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            # Check if the shape is an image
            if hasattr(shape, "image"):
                image = shape.image  # Access the image blob from the shape
                image_bytes = image.blob  # Get the raw image bytes

                # Convert raw bytes to a PIL Image using the utility function
                pil_image = convert_bytes_to_pil_image(image_bytes)

                # Preprocess the image before OCR
                preprocessed_image = preprocess_image(pil_image)

                # Apply OCR to the preprocessed image
                text += pytesseract.image_to_string(preprocessed_image) + ' '
    return text


def extract_image_text_docx(docx_path):
    """Extract images from a Word document and apply OCR with preprocessing."""
    doc = Document(docx_path)
    text = ""
    for rel in doc.part.rels.values():
        # Check if the relationship target is an image
        if "image" in rel.target_ref:
            image_data = rel.target_part.blob  # Get the raw image data

            # Convert raw bytes to a PIL Image using the utility function
            pil_image = convert_bytes_to_pil_image(image_data)

            # Preprocess the image before OCR
            preprocessed_image = preprocess_image(pil_image)

            # Apply OCR to the preprocessed image
            text += pytesseract.image_to_string(preprocessed_image) + ' '
    return text
